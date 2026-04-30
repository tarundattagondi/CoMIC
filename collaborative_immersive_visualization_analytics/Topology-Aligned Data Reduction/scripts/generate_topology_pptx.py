"""Generate a presentation-style PPTX on how topology reduction is used in CIVA.

Output: docs/CIVA_Topology_Reduction_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "CIVA_Topology_Reduction_Presentation.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
GRAY = RGBColor(0x44, 0x4B, 0x52)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bar(slide, color: RGBColor = ACCENT) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.35))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NAVY
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        r = sp.add_run()
        r.text = subtitle
        r.font.size = Pt(14)
        r.font.italic = True
        r.font.color.rgb = GRAY


def add_bullets(slide, bullets: list[str], top: float = 1.9, left: float = 0.6, width: float = 12.1, height: float = 5.3, size: int = 18) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(size)
        run.font.color.rgb = GRAY


def add_textbox(slide, text: str, left: float, top: float, width: float, height: float, size: int = 14, color: RGBColor = GRAY, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_chip(slide, text: str, left: float, top: float, width: float, height: float, fill: RGBColor, fg: RGBColor = WHITE, size: int = 14) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(2.25)


def add_footer(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = GRAY


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    return slide


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----
    s = blank_slide(prs)
    set_bg(s, NAVY)
    add_textbox(s, "Topology-Aligned Data Reduction", 0.6, 2.3, 12.1, 1.0, size=40, color=WHITE, bold=True)
    add_textbox(s, "How CIVA Reduction Uses TTK in an Immersive Analytics Pipeline", 0.6, 3.3, 12.1, 0.7, size=22, color=WHITE)
    add_textbox(
        s,
        "Sai Kiran Annam  •  Tarun Datta Gondi  •  Fahim Asrad Nafis\nMobile & Immersive Computing — George Mason University",
        0.6, 5.5, 12.1, 1.0, size=16, color=WHITE,
    )
    add_bar(s, ACCENT)

    # ---- Slide 2: The Problem ----
    s = blank_slide(prs)
    add_title(s, "The Problem", "Full-resolution volumes do not fit inside a WebXR frame budget")
    add_bullets(
        s,
        [
            "Scientific volumes are huge: 256³ ≈ 16.8 million voxels for a single scalar field.",
            "Immersive rendering demands 72–90 fps; any stall breaks presence and comfort.",
            "Static, offline downsampling cannot adapt to a user's moving focus or task intent.",
            "Uniform decimation blurs small but meaningful features — exactly what analysts care about.",
            "Goal: make reduction interactive, structure-preserving, reversible, and measurable.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 3: Why Topology ----
    s = blank_slide(prs)
    add_title(s, "Why Topology-Aligned Reduction?", "Persistence measures how 'important' a feature is")
    add_bullets(
        s,
        [
            "Every critical point in a scalar field has a persistence value — its lifespan in the filtration.",
            "High persistence = large, salient structure.  Low persistence = noise / insignificant bump.",
            "Thresholding by persistence removes noise without flattening the real features.",
            "Contrast: uniform shrink discards detail everywhere equally, regardless of significance.",
            "Topology-aligned reduction preserves what analysts actually need to see.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 4: Tools ----
    s = blank_slide(prs)
    add_title(s, "What We Use: Topology ToolKit (TTK)", "Plus VTK for I/O and the web-native frontend")
    add_bullets(
        s,
        [
            "TTK (Topology ToolKit) — persistence diagrams & topological simplification filters.",
            "Key filters used: ttkPersistenceDiagram, ttkTopologicalSimplification, ttkTopologicalSimplificationByPersistence.",
            "VTK 9.3 — reads/writes .vti, runs auxiliary filters (vtkThreshold, vtkProbeFilter, vtkImageShrink3D).",
            "FastAPI + uvicorn — exposes reduction as an HTTP endpoint the frontend can call at runtime.",
            "VTK.js + React + TypeScript — browser volume rendering and WebXR session.",
            "Docker Compose — Linux-container TTK build (works on Apple Silicon where conda-forge has no arm64 build).",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 5: Pipeline Diagram ----
    s = blank_slide(prs)
    add_title(s, "How Topology Reduction Runs", "Frontend interaction → backend TTK pipeline → reduced VTI → scene")

    # Row of chips - pipeline
    y = 2.5
    h = 0.9
    x = 0.45
    w = 1.85
    gap = 0.2

    chips = [
        ("User action\n(LOD / persistence)", ACCENT),
        ("FSM / Store\nupdate", ACCENT),
        ("POST /api/reduce", NAVY),
        ("ttkPersistence\nDiagram", NAVY),
        ("vtkThreshold\nby Persistence", NAVY),
        ("ttkTopological\nSimplification", NAVY),
        ("vtkImageShrink3D\n(LOD tier)", ACCENT),
    ]
    positions = []
    for i, (label, color) in enumerate(chips):
        lx = x + i * (w + gap)
        add_chip(s, label, lx, y, w, h, color, WHITE, 12)
        positions.append((lx, y, w, h))

    # arrows between chips
    for i in range(len(positions) - 1):
        lx, ly, lw, lh = positions[i]
        nx, ny, _, _ = positions[i + 1]
        add_arrow(s, lx + lw, ly + lh / 2, nx, ny + lh / 2)

    # final step box below
    add_chip(s, "Reduced VTI + X-Reduction-Metadata → VTK.js renderer → render loop → metrics", 0.8, 4.1, 11.7, 0.8, ACCENT, WHITE, 14)
    add_arrow(s, 6.6, 3.4, 6.6, 4.1)

    add_textbox(
        s,
        "Every arrow is a real module in the repo (backend/reduce.py, backend/main.py, src/data/data-loader.ts, src/state/reduction-fsm.ts, src/metrics/*).",
        0.6, 5.2, 12.1, 0.6, size=13, color=GRAY, bold=False,
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 6: The TTK pipeline inside the backend ----
    s = blank_slide(prs)
    add_title(s, "Inside the Backend: backend/reduce.py", "Three explicit steps + a compatibility fallback")
    add_bullets(
        s,
        [
            "1. Persistence diagram — ttkPersistenceDiagram consumes the .vti's active scalar (e.g. ImageScalars) and emits pairs with a Persistence attribute.",
            "2. Threshold by persistence — vtkThreshold selects only pairs with Persistence ≥ mapped_threshold, producing a set of simplification constraints.",
            "3. Topological simplification — ttkTopologicalSimplification takes the domain + constraints and returns a simplified scalar field (noise collapsed, real features kept).",
            "Fallback — if the explicit two-port pipeline is unavailable in the installed TTK wrapper, we call ttkTopologicalSimplificationByPersistence with SetPersistenceThreshold directly.",
            "Normalize output — if the result is not vtkImageData, vtkProbeFilter resamples it back onto the source grid so the VTI writer can emit it.",
            "Spatial LOD tier — vtkImageShrink3D with factor {full:1, high:1, medium:2, low:4} scales resolution for frame-rate control.",
        ],
        size=16,
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 7: Persistence slider → absolute TTK units ----
    s = blank_slide(prs)
    add_title(s, "From UI Slider to TTK Units", "Frontend sends [0,1]; backend scales it to the dataset's actual persistence range")
    add_bullets(
        s,
        [
            "UI slider range: normalized 0.0 (no simplification) → 1.0 (most aggressive).",
            "Default per LOD level when the user doesn't set one: full=0.0, high=0.0, medium=0.5, low=0.9.",
            "Backend estimates max persistence by scanning the Persistence array from ttkPersistenceDiagram output.",
            "Mapped absolute threshold = normalized_value × estimated_max_persistence.",
            "Raw absolute values > 1 are passed through unchanged (advanced / scripted use).",
            "This keeps the slider meaningful across datasets with very different scalar ranges.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 8: LOD tiers (voxels) ----
    s = blank_slide(prs)
    add_title(s, "LOD Tiers: Spatial Decimation on Top of TTK", "Topology cleans the field; LOD scales resolution to the frame budget")

    # Table-like chips
    headers = ["Level", "Shrink factor", "Voxels", "Use case"]
    rows = [
        ("Full", "1×", "16,777,216 (256³)", "Static reference / screenshots"),
        ("High", "1×", "16,777,216 (256³)", "Default desktop exploration"),
        ("Medium", "2×", "2,097,152 (128³)", "Immersive headset + ROI work"),
        ("Low", "4×", "262,144 (64³)", "Overview fly-through / auto-LOD far camera"),
    ]
    col_w = [1.5, 2.0, 3.0, 5.6]
    col_x = [0.6]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + 0.1)

    top = 2.0
    for i, h in enumerate(headers):
        add_chip(s, h, col_x[i], top, col_w[i], 0.55, NAVY, WHITE, 14)
    for r_idx, row in enumerate(rows):
        ry = top + 0.65 + r_idx * 0.75
        for c_idx, cell in enumerate(row):
            add_chip(s, cell, col_x[c_idx], ry, col_w[c_idx], 0.65, LIGHT, NAVY, 13)

    add_textbox(
        s,
        "Auto-LOD policy: camera distance thresholds trigger tier switches. Enabling the ROI boosts the tier upward inside the ROI window.",
        0.6, 6.3, 12.1, 0.6, size=13, color=GRAY,
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 9: Frontend surface ----
    s = blank_slide(prs)
    add_title(s, "Frontend Surface: What the User Touches", "Every control writes to the FSM and re-derives reduction phase")
    add_bullets(
        s,
        [
            "Persistence slider — normalized [0,1], sent to backend on release.",
            "LOD selector — Full / High / Medium / Low; re-requests reduced VTI.",
            "Auto-LOD toggle — camera-distance-driven tier transitions.",
            "Slice plane — choose dim (X/Y/Z) + offset; half-space clipping, no reload needed.",
            "Contextual dimming — fade non-focused volume to emphasize features.",
            "Isosurface / threshold-region — feature operators on the currently loaded scalar.",
            "ROI sphere — toggle + radius; triggers LOD boost policy inside the window.",
            "Scalar rollback + global reset — reversible runtime state, never destructive.",
        ],
        size=16,
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 10: Reduction FSM ----
    s = blank_slide(prs)
    add_title(s, "Reduction as an Explicit FSM", "Phase is derived from context, not hand-set — no UI drift")

    states = [
        ("idle", 1.0, 3.3),
        ("base_volume", 3.3, 3.3),
        ("lod_switched", 5.8, 2.1),
        ("feature_focus", 5.8, 4.5),
        ("roi_refined", 9.3, 3.3),
    ]
    for name, x, y in states:
        add_chip(s, name, x, y, 2.1, 0.8, ACCENT, WHITE, 16)

    # arrows
    add_arrow(s, 3.1, 3.7, 3.3, 3.7)          # idle → base
    add_arrow(s, 5.4, 3.5, 5.8, 2.5)          # base → lod
    add_arrow(s, 5.4, 3.9, 5.8, 4.7)          # base → feature
    add_arrow(s, 7.9, 2.5, 9.3, 3.5)          # lod → roi
    add_arrow(s, 7.9, 4.7, 9.3, 3.9)          # feature → roi
    add_arrow(s, 5.4, 3.7, 9.3, 3.7)          # base → roi

    add_textbox(
        s,
        "Every transition is logged with a timestamp. The exportable session JSON replays the full reduction trajectory.",
        0.6, 5.9, 12.1, 0.6, size=14, color=GRAY,
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 11: Observability ----
    s = blank_slide(prs)
    add_title(s, "Observability Around Every Reduction", "Numbers are written on action, not guessed afterward")
    add_bullets(
        s,
        [
            "Frame metrics — FPS mean / min / max, frame-time mean & p95 (rolling window).",
            "Action-to-stable-render latency — per kind: LOD, threshold, feature toggle, ROI, volume load.",
            "Reduction metadata — X-Reduction-Metadata header per /api/reduce: usedTTK, reductionMode, mapped threshold, output dims/spacing/origin.",
            "Memory hint — performance.memory snapshot when the browser exposes it.",
            "Session event log — append-only, timestamped, includes interaction kind + payload.",
            "One-click JSON export — events + metric snapshot. Reproducibility is native, not bolted on.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 12: Trade-offs ----
    s = blank_slide(prs)
    add_title(s, "Trade-offs We Can Now Measure", "The dashboard lets us move along each axis deliberately")
    add_bullets(
        s,
        [
            "Performance ↔ Detail — lower LOD lifts FPS but can hide small structures.",
            "Persistence strength ↔ Fidelity — higher threshold kills noise but may erase subtle real features.",
            "Feature emphasis ↔ Global context — slicing/dimming focuses the eye but removes surroundings.",
            "Automation ↔ Predictability — auto-LOD adapts to distance but can surprise the user if thresholds are poor.",
            "All four axes are tied to logged events → analyzable, not anecdotal.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 13: What this gives us ----
    s = blank_slide(prs)
    add_title(s, "What This Architecture Gives Us", "Why topology-aligned reduction is worth the integration cost")
    add_bullets(
        s,
        [
            "Reduction is a live system behavior, not a frozen preprocessing artifact.",
            "Structure is preserved: noise is removed but features with high persistence survive.",
            "Every interaction is reversible — scalar rollback, ROI disable, global reset, no lost exploration.",
            "Deterministic state: phase is recomputed from context; the same actions always land in the same phase.",
            "Session-level reproducibility: export a run, replay its decisions, compare metrics across runs.",
            "Works in both desktop and WebXR modes with the same core pipeline.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 14: Limitations + Future ----
    s = blank_slide(prs)
    add_title(s, "Limitations & Where We Go Next", "Honest view of what remains on the table")
    add_bullets(
        s,
        [
            "Currently validated primarily on one canonical dataset (ctBones 256³).",
            "ROI is a visualization + policy primitive, not yet full localized high-res subvolume replacement.",
            "TTK on osx-arm64 still requires Docker; native conda-forge build does not exist for Apple Silicon.",
            "Future: ROI-local high-resolution streaming with composited rendering.",
            "Future: persistence-informed isosurface defaults + threshold segmentation operators.",
            "Future: scripted scenario runners for automated cross-browser/device benchmark traces.",
            "Future: adaptive reduction policies that jointly optimize FPS and feature-preservation.",
        ],
    )
    add_footer(s, "CIVA Reduction — Topology Reduction Deck")

    # ---- Slide 15: Summary ----
    s = blank_slide(prs)
    set_bg(s, NAVY)
    add_bar(s, ACCENT)
    add_textbox(s, "Summary", 0.6, 1.0, 12.1, 1.0, size=36, color=WHITE, bold=True)
    add_textbox(
        s,
        "Topology reduction in CIVA is a runtime pipeline:\n"
        "persistence diagram  →  threshold by persistence  →  topological simplification  →  LOD shrink.\n\n"
        "Driven by live UI controls, gated by an explicit FSM, measured end-to-end, exportable as JSON,\n"
        "and served to a VTK.js / WebXR frontend that can reverse any step without reloading the session.",
        0.7, 2.3, 12.0, 4.2, size=20, color=WHITE,
    )
    add_textbox(s, "Questions?", 0.6, 6.2, 12.1, 0.8, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    return prs


def main() -> None:
    prs = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote: {OUT}")


if __name__ == "__main__":
    main()
